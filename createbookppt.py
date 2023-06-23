import sqlite3
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
from bs4 import BeautifulSoup


# Function to convert PNG image to JPEG with white background
def convert_image(image_path, converted_image_path):
    image = Image.open(image_path)
    image = image.convert("RGBA")
    background = Image.new("RGB", image.size, (255, 255, 255))
    background.paste(image, mask=image.split()[3])
    background.save(converted_image_path, "JPEG", quality=100)


# Connect to the SQLite database
conn = sqlite3.connect('e:/qeraat/data_v15.db')
cursor = conn.cursor()

# Execute the query to retrieve data from both tables
query = '''
SELECT mj.text, ms.page_number
FROM book_jlalin AS mj
JOIN mosshf_shmrly AS ms ON mj.aya_index = ms.aya_index
'''
cursor.execute(query)
data = cursor.fetchall()

# Create a new PowerPoint presentation
prs = Presentation()

# Initialize variables for tracking current page number and concatenated text
current_page = None
concatenated_text = ""

# Iterate over the retrieved data
for text, page_number in data:
    # Check if the page has changed
    if current_page is None or current_page != page_number:
        # Add the previous slide (except for the first iteration)
        if current_page is not None:
            # Create a new slide with a title and content layout
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)

            # Add the image to the slide
            image_path = f'E:/Qeraat/pages/{current_page}.png'  # Replace with the actual path to the folder and image file extension
            converted_image_path = f'E:/Qeraat/pages/{current_page}.jpg'  # Replace with the path and filename for the converted image
            convert_image(image_path, converted_image_path)
            left = Inches(0.5)
            top = Inches(1.5)
            width = Inches(3)
            height = Inches(4)
            slide.shapes.add_picture(converted_image_path, left, top, width, height)

            # Add the concatenated text to the slide
            text_box = slide.shapes.add_textbox(left + width + Inches(0.5), top, Inches(6 - 0.5 - 3 - 0.5), height)
            text_frame = text_box.text_frame
            text_frame.word_wrap = True
            p = text_frame.add_paragraph()
            p.text = concatenated_text

        # Update the current page number and reset the concatenated text
        current_page = page_number
        concatenated_text = ""

    # Process the HTML text and concatenate it
    soup = BeautifulSoup(text, 'html.parser')
    concatenated_text += soup.get_text() + " "  # Modify the separator as needed

# Add the last slide to the presentation
slide_layout = prs.slide_layouts[1]  # Title and Content layout
slide = prs.slides.add_slide(slide_layout)

# Add the image to the slide
image_path = f'E:/Qeraat/pages/{current_page}.png'  # Replace with the actual path to the folder and image file extension
converted_image_path = f'E:/Qeraat/pages/{current_page}.jpg'  # Replace with the path and filename for the converted image
convert_image(image_path, converted_image_path)
left = Inches(0.5)
top = Inches(1.5)
width = Inches(3)
height = Inches(4)
slide.shapes.add_picture(converted_image_path, left, top, width, height)

# Add the concatenated text to the slide
text_box = slide.shapes.add_textbox(left + width + Inches(0.5), top, Inches(6 - 0.5 - 3 - 0.5), height)
text_frame = text_box.text_frame
text_frame.word_wrap = True
p = text_frame.add_paragraph()
p.text = concatenated_text

# Save the PowerPoint presentation
prs.save('output.pptx')

# Close the database connection
conn.close()
